"""
Renders the built .pptx to PNGs so the layout can be checked without Office.

    python -m pip install pillow
    python docs/sales-deck/preview_deck.py

This reads the saved presentation back through python-pptx and draws every
shape with Pillow, using the same Segoe UI files PowerPoint will use. It is an
approximation — PowerPoint does its own kerning and hyphenation — but it is
faithful enough to catch the two failures that matter: text that overflows its
box, and blocks that collide.
"""

from __future__ import annotations

from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu

import metrics

HERE = Path(__file__).parent
DECK = HERE / "MyCampusView-Sales-Deck.pptx"
OUT = HERE / "preview"
DPI = 150

FONTS = {
    "Segoe UI": "C:/Windows/Fonts/segoeui.ttf",
    "Segoe UI Semibold": "C:/Windows/Fonts/seguisb.ttf",
    "Segoe UI Light": "C:/Windows/Fonts/segoeuil.ttf",
    "Segoe UI Semilight": "C:/Windows/Fonts/segoeuisl.ttf",
}
_cache: dict = {}


def font(name, size_pt):
    key = (name, round(size_pt, 1))
    if key not in _cache:
        path = FONTS.get(name, FONTS["Segoe UI"])
        _cache[key] = ImageFont.truetype(path, int(round(size_pt * DPI / 72)))
    return _cache[key]


def px(emu):
    return int(round(Emu(int(emu)).inches * DPI))


def solid_fill(shape):
    """(r, g, b) or None. Reads the alpha the build script may have written."""
    try:
        f = shape.fill
        if f.type is None or f.type == 5:  # background / none
            return None
        rgb = f.fore_color.rgb
        alpha = 1.0
        el = f.fore_color._xFill.find(qn("a:srgbClr"))
        if el is not None:
            a = el.find(qn("a:alpha"))
            if a is not None:
                alpha = int(a.get("val")) / 100000
        return (rgb[0], rgb[1], rgb[2]), alpha
    except Exception:
        return None


def gradient_fill(shape):
    """(stops, angle_degrees) for the alpha gradients the build script writes."""
    spPr = getattr(shape._element, "spPr", None)
    if spPr is None:
        return None
    grad = spPr.find(qn("a:gradFill"))
    if grad is None:
        return None
    stops = []
    for gs in grad.findall(qn("a:gsLst") + "/" + qn("a:gs")) or grad.find(qn("a:gsLst")):
        pos = int(gs.get("pos", "0")) / 100000
        clr = gs.find(qn("a:srgbClr"))
        if clr is None:
            continue
        v = clr.get("val")
        rgb = (int(v[0:2], 16), int(v[2:4], 16), int(v[4:6], 16))
        a = clr.find(qn("a:alpha"))
        alpha = int(a.get("val")) / 100000 if a is not None else 1.0
        stops.append((pos, rgb, alpha))
    lin = grad.find(qn("a:lin"))
    ang = int(lin.get("ang", "0")) / 60000 if lin is not None else 0
    return (sorted(stops), ang) if stops else None


def render_gradient(img, shape, stops, angle_deg):
    import numpy as np

    x, y = px(shape.left), px(shape.top)
    w, h = max(1, px(shape.width)), max(1, px(shape.height))
    yy, xx = np.mgrid[0:h, 0:w]
    rad = np.deg2rad(angle_deg)
    # Project onto the gradient direction, normalised over the shape's extent.
    proj = xx * np.cos(rad) + yy * np.sin(rad)
    proj -= proj.min()
    span = proj.max() or 1
    t = proj / span

    pos = np.array([s[0] for s in stops])
    alpha = np.array([s[2] for s in stops])
    cols = np.array([s[1] for s in stops], dtype=float)
    a = np.interp(t, pos, alpha)
    r = np.interp(t, pos, cols[:, 0])
    g = np.interp(t, pos, cols[:, 1])
    b = np.interp(t, pos, cols[:, 2])

    layer = np.dstack([r, g, b, a * 255]).astype("uint8")
    img.alpha_composite(Image.fromarray(layer, "RGBA"), (x, y))


def line_colour(shape):
    try:
        if shape.line.fill.type in (None, 5):
            return None
        c = shape.line.color.rgb
        return (c[0], c[1], c[2])
    except Exception:
        return None


def wrap(body, name, size_pt, max_px):
    """
    Wrap through `metrics`, so the preview breaks lines exactly where the
    builder assumed they would break. Any divergence between the two would
    make the preview useless for spotting collisions.
    """
    return metrics.wrap(body, name, size_pt, max_px * 72 / DPI)


def render_shape(img, draw, shape, overflow):
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return render_picture(img, shape)

    x, y, w, h = px(shape.left), px(shape.top), px(shape.width), px(shape.height)

    grad = gradient_fill(shape)
    if grad:
        render_gradient(img, shape, *grad)
        if not shape.has_text_frame:
            return
        render_text(img, draw, shape, x, y, w, h, overflow)
        return

    fill = solid_fill(shape)
    stroke = line_colour(shape)
    if fill or stroke:
        rgba = None
        if fill:
            (r, g, b), alpha = fill
            rgba = (r, g, b, int(alpha * 255))
        layer = Image.new("RGBA", img.size, (0, 0, 0, 0))
        ld = ImageDraw.Draw(layer)
        name = (shape.shape_type or "")
        try:
            auto = shape.auto_shape_type
        except Exception:
            auto = None
        box = [x, y, x + w, y + h]
        if auto is not None and "OVAL" in str(auto):
            ld.ellipse(box, fill=rgba, outline=stroke, width=1)
        elif auto is not None and "ROUNDED" in str(auto):
            radius = int(min(w, h) * 0.5 * (shape.adjustments[0] if len(shape.adjustments) else 0.16))
            ld.rounded_rectangle(box, radius=max(1, radius), fill=rgba, outline=stroke, width=1)
        else:
            ld.rectangle(box, fill=rgba, outline=stroke, width=1)
        img.alpha_composite(layer)

    if not shape.has_text_frame:
        return
    render_text(img, draw, shape, x, y, w, h, overflow)


def render_text(img, draw, shape, x, y, w, h, overflow):
    tf = shape.text_frame
    ml = px(tf.margin_left or 0)
    mr = px(tf.margin_right or 0)
    mt = px(tf.margin_top or 0)
    max_w = max(10, w - ml - mr)

    # Vertical anchoring: measure first, then place.
    blocks = []
    for p in tf.paragraphs:
        runs = [r for r in p.runs if r.text]
        if not runs:
            blocks.append((None, [], 0))
            continue
        size = runs[0].font.size.pt if runs[0].font.size else 12
        name = runs[0].font.name or "Segoe UI"
        fnt = font(name, size)
        body = "".join(r.text for r in runs)
        lines = wrap(body, name, size, max_w)
        # The builder writes exact point spacing, which is what PowerPoint
        # honours too; a float would be a multiple of the font's own leading.
        ls = p.line_spacing
        lh = (ls.pt if hasattr(ls, "pt") else (ls or 1.2) * size) * DPI / 72
        after = (p.space_after.pt * DPI / 72) if p.space_after else 0
        colour = runs[0].font.color.rgb if runs[0].font.color and runs[0].font.color.type is not None else None
        rgb = (colour[0], colour[1], colour[2]) if colour else (0, 0, 0)
        blocks.append(((fnt, rgb, p.alignment, lh), lines, len(lines) * lh + after))

    total = sum(b[2] for b in blocks)
    anchor = str(tf.vertical_anchor or "")
    if "MIDDLE" in anchor:
        cy = y + (h - total) / 2
    elif "BOTTOM" in anchor:
        cy = y + h - total
    else:
        cy = y + mt

    top0 = cy
    for meta, lines, _ in blocks:
        if meta is None:
            cy += 6
            continue
        fnt, rgb, align, lh = meta
        for line in lines:
            tw = draw.textlength(line, font=fnt)
            if align is not None and "CENTER" in str(align):
                tx = x + ml + (max_w - tw) / 2
            elif align is not None and "RIGHT" in str(align):
                tx = x + ml + max_w - tw
            else:
                tx = x + ml
            draw.text((tx, cy), line, font=fnt, fill=rgb)
            cy += lh

    # Report anything that ran past the bottom of its own box by a real margin.
    if cy - top0 > h + 8 and shape.width > Emu(200000):
        overflow.append((shape.shape_id, round((cy - top0 - h) / DPI, 2),
                         (shape.text_frame.text or "")[:58].replace("\n", " ")))


def render_picture(img, shape):
    blob = shape.image.blob
    from io import BytesIO
    src = Image.open(BytesIO(blob)).convert("RGBA")
    sw, sh = src.size
    l = shape.crop_left or 0
    t = shape.crop_top or 0
    r = shape.crop_right or 0
    b = shape.crop_bottom or 0
    box = (int(sw * l), int(sh * t), int(sw * (1 - r)), int(sh * (1 - b)))
    if box[2] > box[0] and box[3] > box[1]:
        src = src.crop(box)
    w, h = max(1, px(shape.width)), max(1, px(shape.height))
    src = src.resize((w, h), Image.LANCZOS)
    img.alpha_composite(src, (px(shape.left), px(shape.top)))


def main():
    OUT.mkdir(exist_ok=True)
    prs = Presentation(str(DECK))
    W, H = px(prs.slide_width), px(prs.slide_height)

    for i, slide in enumerate(prs.slides, start=1):
        img = Image.new("RGBA", (W, H), (255, 255, 255, 255))
        draw = ImageDraw.Draw(img)
        overflow = []
        for shape in slide.shapes:
            try:
                render_shape(img, draw, shape, overflow)
            except Exception as err:  # a preview must never stop at one shape
                print(f"  slide {i}: skipped {shape.shape_id} — {err}")
        img.convert("RGB").save(OUT / f"slide-{i:02d}.png", quality=95)

        # A text box sized to its own content never "overflows", so the check
        # above cannot catch a block that simply runs off the foot of the
        # slide. This one can.
        for shape in slide.shapes:
            if not shape.has_text_frame or not (shape.text_frame.text or "").strip():
                continue
            bottom = shape.top + shape.height
            if bottom > prs.slide_height - Emu(180000):
                sample = " ".join(shape.text_frame.text.split())[:52]
                over = round((bottom - prs.slide_height) / 914400, 2)
                print(f"slide {i}: runs past the foot by {over}in  “{sample}…”")

        if overflow:
            print(f"slide {i}: {len(overflow)} box(es) overflowing")
            for sid, over, sample in overflow:
                print(f"    +{over}in  “{sample}…”")
    print(f"\n{OUT}")


if __name__ == "__main__":
    main()
