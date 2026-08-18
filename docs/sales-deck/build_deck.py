"""
Builds the MyCampusView sales deck.

    python -m pip install python-pptx pillow
    python docs/sales-deck/build_deck.py
    python docs/sales-deck/preview_deck.py     # renders PNGs so it can be checked

Design follows the marketing site rather than a template: the alternation of
black and paper grounds the homepage uses as chapters, the same editorial
accents, one display face set large, and a great deal of air. Copy lives in
`deck_content.py`; nothing here invents a claim.

Every text block is measured before the next is placed (see `metrics.py`) and
every line height is written as an exact point value, so the preview and
PowerPoint agree about where things land.

Screenshots in `screenshots/` are real captures of the running application
against the seeded demonstration school. Re-capture them rather than editing
them — a retouched screenshot is a claim we cannot stand behind.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches, Pt

import deck_content as C
import metrics

HERE = Path(__file__).parent
SHOTS = HERE / "screenshots"
BRAND = HERE.parent.parent / "public" / "brand"
OUT = HERE / "MyCampusView-Sales-Deck.pptx"

# ------------------------------------------------------------------ palette
# Site tokens from src/styles/site.css, plus the violet the application renders
# its primary actions in.
INK = RGBColor(0x0A, 0x10, 0x24)
INK_DEEP = RGBColor(0x06, 0x0A, 0x18)
PAPER = RGBColor(0xFB, 0xFA, 0xF8)
PAPER_2 = RGBColor(0xF4, 0xF3, 0xEF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK_SOFT = RGBColor(0x53, 0x5B, 0x6E)
INK_FAINT = RGBColor(0x8A, 0x92, 0xA3)
ON_DARK = RGBColor(0xF6, 0xF5, 0xF2)
ON_DARK_SOFT = RGBColor(0xAD, 0xB3, 0xC6)
ON_DARK_FAINT = RGBColor(0x71, 0x79, 0x8F)
RULE = RGBColor(0xDF, 0xE3, 0xEC)
RULE_DARK = RGBColor(0x28, 0x2E, 0x45)

VIOLET = RGBColor(0x6C, 0x4B, 0xF4)
VIOLET_SOFT = RGBColor(0xB5, 0x7F, 0xF0)
AMBER = RGBColor(0xF0, 0xA6, 0x3C)
ROSE = RGBColor(0xEF, 0x5B, 0xA8)
SKY = RGBColor(0x4E, 0xA3, 0xF5)
MINT = RGBColor(0x2F, 0xB5, 0x73)
ACCENTS = [SKY, ROSE, AMBER, MINT, VIOLET_SOFT]

DISPLAY = "Segoe UI Semibold"
BODY = "Segoe UI"

W, H = Inches(13.333), Inches(7.5)
M = Inches(0.86)                      # outer margin
CONTENT = W - M * 2                   # 11.613in of usable width


def IN(emu) -> float:
    return emu / 914400


# ------------------------------------------------------------- raw XML bits

def _shadow(shape, blur=28, dist=10, alpha=0.13):
    """A soft drop shadow. python-pptx exposes no API for one."""
    spPr = shape._element.spPr
    for old in spPr.findall(qn("a:effectLst")):
        spPr.remove(old)
    spPr.append(parse_xml(
        f'<a:effectLst {nsdecls("a")}>'
        f'<a:outerShdw blurRad="{int(blur * 12700)}" dist="{int(dist * 12700)}" '
        f'dir="5400000" rotWithShape="0">'
        f'<a:srgbClr val="0A1024"><a:alpha val="{int(alpha * 100000)}"/></a:srgbClr>'
        f'</a:outerShdw></a:effectLst>'
    ))


def _no_shadow(shape):
    shape.shadow.inherit = False


def _tracking(run, hundredths):
    run.font._rPr.set("spc", str(int(hundredths)))


def _send_to_back(slide, shape, index=2):
    tree = slide.shapes._spTree
    tree.remove(shape._element)
    tree.insert(index, shape._element)


def _put_behind(slide, shape, other):
    tree = slide.shapes._spTree
    tree.remove(shape._element)
    tree.insert(list(tree).index(other._element), shape._element)


# ------------------------------------------------------------------ shapes

def rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, radius=None):
    s = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if radius:
        s.adjustments[0] = radius
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    _no_shadow(s)
    return s


def oval(slide, x, y, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    _no_shadow(s)
    return s


def hairline(slide, x, y, w, colour, weight=0.9):
    return rect(slide, x, y, w, Emu(int(Pt(weight))), fill=colour)


CORNER_ANGLE = {"tr": 8100000, "br": 13500000, "tl": 2700000, "bl": 18900000}


def wash(slide, colour, corner, peak=0.34):
    """
    Full-bleed glow anchored to a corner.

    A gradient rectangle that stops inside the frame leaves a hard seam on the
    edges the gradient does not run along, so it always bleeds off all four.
    """
    s = rect(slide, 0, 0, W, H)
    hexed = f"{colour[0]:02X}{colour[1]:02X}{colour[2]:02X}"
    spPr = s._element.spPr
    for tag in ("a:solidFill", "a:noFill", "a:gradFill"):
        for old in spPr.findall(qn(tag)):
            spPr.remove(old)
    grad = parse_xml(
        f'<a:gradFill {nsdecls("a")} rotWithShape="1"><a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{hexed}">'
        f'<a:alpha val="{int(peak * 100000)}"/></a:srgbClr></a:gs>'
        f'<a:gs pos="45000"><a:srgbClr val="{hexed}">'
        f'<a:alpha val="{int(peak * 34000)}"/></a:srgbClr></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{hexed}">'
        f'<a:alpha val="0"/></a:srgbClr></a:gs>'
        f'</a:gsLst><a:lin ang="{CORNER_ANGLE[corner]}" scaled="0"/></a:gradFill>'
    )
    ln = spPr.find(qn("a:ln"))
    ln.addprevious(grad) if ln is not None else spPr.append(grad)
    return s


# -------------------------------------------------------------------- text

def block(slide, x, y, w, body, *, size=12, colour=INK, font=BODY, leading=1.35,
          align=PP_ALIGN.LEFT, spc=0, anchor=MSO_ANCHOR.TOP):
    """
    One measured paragraph. Returns the height it occupies, in EMU.

    Line spacing is written as exact points, so the height computed here is the
    height PowerPoint will use.
    """
    h_in = metrics.height_in(body, font, size, IN(w), leading, spc)
    box = slide.shapes.add_textbox(x, y, w, Inches(h_in) + Inches(0.02))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor

    for i, line in enumerate(body.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = Pt(size * leading)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = colour
        r.font.name = font
        if spc:
            _tracking(r, spc)
    return Inches(h_in)


def eyebrow(slide, x, y, label, colour, size=10.5):
    return block(slide, x, y, Inches(7), label.upper(), size=size, colour=colour,
                 font=DISPLAY, leading=1.3, spc=180)


def headline(slide, x, y, w, body, colour, size=38, leading=1.08):
    return block(slide, x, y, w, body, size=size, colour=colour, font=DISPLAY,
                 leading=leading)


# --------------------------------------------------------------- components

def running_mark(slide, dark: bool):
    """The lockup, small, as a running header."""
    lockup(slide, M, Inches(0.46), Inches(0.19), dark=dark)


def page_number(slide, n, dark: bool):
    block(slide, W - M - Inches(1.2), H - Inches(0.6), Inches(1.2), f"{n:02d}",
          size=9.5, colour=ON_DARK_FAINT if dark else INK_FAINT, spc=120,
          align=PP_ALIGN.RIGHT)


def ground(slide, dark: bool, tone=None):
    bg = rect(slide, 0, 0, W, H,
              fill=tone if tone is not None else (INK if dark else PAPER))
    _send_to_back(slide, bg)
    return bg


def lockup(slide, x, y, height, *, dark: bool):
    """
    The MyCampusView lockup, placed straight onto the slide.

    Both cuts are transparent and cropped to the artwork, so the slide's own
    ground is the logo's background — there is no plate, and nothing to align
    but the logo itself. `dark` picks the cream wordmark the artwork is drawn
    with; the light cut sets those glyphs in the brand navy so they survive a
    paper ground.
    """
    src = BRAND / ("mycampusview-lockup.png" if dark else "mycampusview-lockup-dark.png")
    if not src.exists():
        return None
    return slide.shapes.add_picture(str(src), x, y, height=height)


def screenshot(slide, name, x, y, w, *, crop=(0, 0, 0, 0), caption=None,
               caption_colour=INK_FAINT, frame=True):
    """
    A capture inside a light frame with a soft shadow.

    `crop` is (left, top, right, bottom) as fractions of the source. The height
    is derived from what survives the crop, so nothing is ever stretched.
    """
    src = SHOTS / name
    if not src.exists():
        raise FileNotFoundError(src)

    pic = slide.shapes.add_picture(str(src), x, y, width=w)
    pic.crop_left, pic.crop_top, pic.crop_right, pic.crop_bottom = crop
    nw, nh = pic.image.size
    kept_w = nw * (1 - crop[0] - crop[2])
    kept_h = nh * (1 - crop[1] - crop[3])
    pic.width = w
    pic.height = Emu(int(w * kept_h / kept_w))

    if frame:
        pad = Inches(0.05)
        f = rect(slide, x - pad, y - pad, pic.width + pad * 2, pic.height + pad * 2,
                 fill=WHITE, line=RULE, radius=0.018)
        _shadow(f, blur=36, dist=13, alpha=0.17)
        _put_behind(slide, f, pic)

    if caption:
        block(slide, x, y + pic.height + Inches(0.17), w, caption, size=9.5,
              colour=caption_colour, leading=1.3)
    return pic


def bullets(slide, x, y, w, items, *, size=11, colour=INK_SOFT, dot=VIOLET,
            gap=Inches(0.17), leading=1.34):
    """Dotted list. Each item is measured, so two-line items never collide."""
    cy = y
    tx = x + Inches(0.2)
    tw = w - Inches(0.2)
    for item in items:
        oval(slide, x, cy + Inches(size * leading / 72 / 2) - Inches(0.028),
             Inches(0.055), Inches(0.055), dot)
        cy += block(slide, tx, cy, tw, item, size=size, colour=colour, leading=leading)
        cy += gap
    return cy


def defs(slide, x, y, w, items, *, size_t=12, size_b=10.5, colour=INK, soft=INK_SOFT,
         rule=RULE, gap=Inches(0.24), leading_t=1.18, leading_b=1.36):
    """Term-and-definition stack under hairlines — the site's own list idiom."""
    cy = y
    for title, body in items:
        hairline(slide, x, cy, w, rule)
        cy += Inches(0.14)
        cy += block(slide, x, cy, w, title, size=size_t, colour=colour, font=DISPLAY,
                    leading=leading_t)
        cy += Inches(0.055)
        cy += block(slide, x, cy, w, body, size=size_b, colour=soft, leading=leading_b)
        cy += gap
    return cy


def _def_height(title, body, w, size_t, size_b, leading_t=1.18, leading_b=1.36):
    return (Inches(0.14)
            + Inches(metrics.height_in(title, DISPLAY, size_t, IN(w), leading_t))
            + Inches(0.055)
            + Inches(metrics.height_in(body, BODY, size_b, IN(w), leading_b)))


def def_grid(slide, x, y, w, items, cols, *, size_t=12, size_b=10, gap_x=Inches(0.5),
             gap_y=Inches(0.34), colour=INK, soft=INK_SOFT, rule=RULE):
    """
    A grid of term/definition cells with per-row heights.

    Rows are measured before anything is drawn, so a two-line title in one cell
    pushes its whole row down rather than overlapping its own body text.
    """
    col_w = Emu(int((w - gap_x * (cols - 1)) / cols))
    heights = [_def_height(t, b, col_w, size_t, size_b) for t, b in items]
    rows = (len(items) + cols - 1) // cols
    row_h = [max(heights[r * cols:(r + 1) * cols]) for r in range(rows)]

    for i, (title, body) in enumerate(items):
        col, row = i % cols, i // cols
        cx = x + col * (col_w + gap_x)
        cy = y + sum(row_h[:row]) + gap_y * row
        hairline(slide, cx, cy, col_w, rule)
        ty = cy + Inches(0.14)
        ty += block(slide, cx, ty, col_w, title, size=size_t, colour=colour,
                    font=DISPLAY, leading=1.18)
        block(slide, cx, ty + Inches(0.055), col_w, body, size=size_b, colour=soft,
              leading=1.36)
    return y + sum(row_h) + gap_y * (rows - 1)


def listing(slide, x, y, w, head, items, accent, *, head_colour=INK,
            item_colour=INK_SOFT, size_h=9.5, size=9.5, leading=1.36,
            gap=Inches(0.075)):
    """A captioned column of short lines, under a coloured rule."""
    hairline(slide, x, y, w, accent, weight=1.8)
    cy = y + Inches(0.16)
    cy += block(slide, x, cy, w, head.upper(), size=size_h, colour=head_colour,
                font=DISPLAY, spc=130, leading=1.25)
    cy += Inches(0.14)
    for item in items:
        cy += block(slide, x, cy, w, item, size=size, colour=item_colour,
                    leading=leading)
        cy += gap
    return cy


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def notes(slide, body):
    slide.notes_slide.notes_text_frame.text = body


# ------------------------------------------------------------------- slides

def slide_cover(prs, d):
    s = new_slide(prs)
    ground(s, dark=True, tone=INK_DEEP)
    wash(s, VIOLET, "tr", peak=0.44)
    wash(s, ROSE, "br", peak=0.15)

    lockup(s, M, Inches(0.82), Inches(0.62), dark=True)

    y = Inches(3.02)
    y += eyebrow(s, M, y, d["eyebrow"], VIOLET_SOFT) + Inches(0.26)
    y += headline(s, M, y, Inches(8.2), d["title"], ON_DARK, size=46, leading=1.06)
    y += Inches(0.34)
    block(s, M, y, Inches(6.9), d["lead"], size=14, colour=ON_DARK_SOFT, leading=1.5)

    hairline(s, M, Inches(6.5), CONTENT, RULE_DARK)
    block(s, M, Inches(6.72), Inches(7.4), d["footer"], size=11.5, colour=ON_DARK_SOFT)
    block(s, M, Inches(6.72), CONTENT, C.TAGLINE, size=11.5, colour=AMBER,
          font=DISPLAY, spc=60, align=PP_ALIGN.RIGHT)

    notes(s, d["notes"])
    return s


def slide_challenge(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False)
    running_mark(s, dark=False)

    y = Inches(1.16)
    y += eyebrow(s, M, y, d["eyebrow"], ROSE) + Inches(0.22)
    y += headline(s, M, y, Inches(6.4), d["title"], INK, size=35)
    y += Inches(0.3)
    block(s, M, y, Inches(5.5), d["lead"], size=13, colour=INK_SOFT, leading=1.5)

    # The kicker sits opposite the lead, set larger — it is the argument.
    block(s, Inches(7.5), Inches(2.36), Inches(5.0), d["kicker"], size=14,
          colour=INK, font=DISPLAY, leading=1.44)

    hairline(s, M, Inches(3.86), CONTENT, RULE)
    def_grid(s, M, Inches(4.24), CONTENT, d["points"], cols=3,
             size_t=12, size_b=9.8, gap_x=Inches(0.52), gap_y=Inches(0.66))

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_connected(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=True)
    wash(s, VIOLET, "br", peak=0.20)
    running_mark(s, dark=True)

    y = Inches(1.2)
    y += eyebrow(s, M, y, d["eyebrow"], AMBER) + Inches(0.22)
    y += headline(s, M, y, Inches(5.0), d["title"], ON_DARK, size=35)
    y += Inches(0.32)
    block(s, M, y, Inches(4.3), d["lead"], size=12.5, colour=ON_DARK_SOFT, leading=1.5)

    hairline(s, M, Inches(5.2), Inches(4.3), RULE_DARK)
    block(s, M, Inches(5.46), Inches(4.3), d["kicker"], size=13.5, colour=ON_DARK,
          font=DISPLAY, leading=1.4)

    col_w = Inches(3.24)
    x1 = Inches(5.66)
    x2 = x1 + col_w + Inches(0.62)

    # Size the card to the taller column rather than to a guess.
    def column_height(items):
        h = Inches(0.16) + Inches(0.28)
        h += Inches(metrics.height_in(d["left_head"], DISPLAY, 9.5, IN(col_w), 1.25, 140))
        for item in items:
            h += Inches(metrics.height_in(item, BODY, 11, IN(col_w), 1.4)) + Inches(0.27)
        return h

    tallest = max(column_height(d["left"]), column_height(d["right"]))
    card = rect(s, x2 - Inches(0.4), Inches(1.16), col_w + Inches(0.8),
                tallest + Inches(0.5), fill=RGBColor(0x12, 0x18, 0x2E), radius=0.045)
    _no_shadow(card)

    for x, head, items, head_colour, item_colour in [
        (x1, d["left_head"], d["left"], ON_DARK_FAINT, ON_DARK_SOFT),
        (x2, d["right_head"], d["right"], MINT, ON_DARK),
    ]:
        cy = Inches(1.52)
        cy += block(s, x, cy, col_w, head.upper(), size=9.5, colour=head_colour,
                    font=DISPLAY, spc=140, leading=1.25)
        cy += Inches(0.16)
        hairline(s, x, cy, col_w, RULE_DARK)
        cy += Inches(0.28)
        for item in items:
            cy += block(s, x, cy, col_w, item, size=11, colour=item_colour, leading=1.4)
            cy += Inches(0.27)

    page_number(s, n, dark=True)
    notes(s, d["notes"])
    return s


def slide_what(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False)
    running_mark(s, dark=False)

    y = Inches(1.16)
    y += eyebrow(s, M, y, d["eyebrow"], VIOLET) + Inches(0.22)
    y += headline(s, M, y, Inches(5.0), d["title"], INK, size=37)
    y += Inches(0.32)
    block(s, M, y, Inches(4.8), d["lead"], size=12.5, colour=INK_SOFT, leading=1.5)

    hairline(s, M, Inches(5.3), Inches(4.8), RULE)
    block(s, M, Inches(5.54), Inches(4.8), d["kicker"], size=10.5, colour=INK_SOFT,
          leading=1.46)

    x = Inches(6.86)
    card_w = W - M - x
    tw = card_w - Inches(1.5)
    cy = Inches(1.18)
    for i, (abbr, name, body) in enumerate(d["products"]):
        body_h = Inches(metrics.height_in(body, BODY, 10.5, IN(tw), 1.4))
        card_h = Inches(0.68) + Inches(0.3) + body_h
        card = rect(s, x, cy, card_w, card_h, fill=WHITE, line=RULE, radius=0.055)
        _shadow(card, blur=22, dist=7, alpha=0.07)

        chip_w, chip_h = Inches(0.64), Inches(0.31)
        rect(s, x + Inches(0.32), cy + Inches(0.34), chip_w, chip_h,
             fill=ACCENTS[i], radius=0.5)
        block(s, x + Inches(0.32), cy + Inches(0.395), chip_w, abbr, size=9,
              colour=INK if ACCENTS[i] is AMBER else WHITE, font=DISPLAY,
              spc=80, align=PP_ALIGN.CENTER)

        tx = x + Inches(1.14)
        block(s, tx, cy + Inches(0.31), tw, name, size=15, colour=INK, font=DISPLAY,
              leading=1.2)
        block(s, tx, cy + Inches(0.7), tw, body, size=10.5, colour=INK_SOFT, leading=1.4)
        cy += card_h + Inches(0.3)

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_value(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False, tone=PAPER_2)
    running_mark(s, dark=False)

    y = Inches(1.16)
    y += eyebrow(s, M, y, d["eyebrow"], MINT) + Inches(0.22)
    y += headline(s, M, y, Inches(6.4), d["title"], INK, size=33)
    y += Inches(0.28)
    block(s, M, y, Inches(5.2), d["lead"], size=12.5, colour=INK_SOFT, leading=1.5)

    # The four figures, as the site sets them: colour discs. Labels share a
    # baseline — staggering the discs read as decoration rather than as data.
    disc, step, dy = Inches(1.12), Inches(1.56), Inches(4.06)
    label_y = dy + disc + Inches(0.2)
    for i, (value, label, note) in enumerate(d["figures"]):
        x = M + i * step
        oval(s, x, dy, disc, disc, ACCENTS[i])
        block(s, x, dy + Inches(0.35), disc, value, size=22 if len(value) < 4 else 18,
              colour=INK, font=DISPLAY, align=PP_ALIGN.CENTER)
        ly = label_y
        ly += block(s, x - Inches(0.14), ly, disc + Inches(0.28), label, size=9.5,
                    colour=INK, font=DISPLAY, align=PP_ALIGN.CENTER, leading=1.25)
        block(s, x - Inches(0.18), ly + Inches(0.06), disc + Inches(0.36), note,
              size=8, colour=INK_SOFT, align=PP_ALIGN.CENTER, leading=1.32)

    x = Inches(7.66)
    inner = W - M - x
    top, pad = Inches(1.16), Inches(0.42)
    body_h = sum(
        (_def_height(t, b, inner, 12, 10) + Inches(0.3)) for t, b in d["points"]
    ) - Inches(0.3)
    card = rect(s, x - Inches(0.46), top, inner + Inches(0.92), body_h + pad * 2,
                fill=WHITE, line=RULE, radius=0.032)
    _shadow(card, blur=26, dist=8, alpha=0.08)
    defs(s, x, top + pad, inner, d["points"], size_t=12, size_b=10, gap=Inches(0.3))

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_platform(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=True)
    wash(s, VIOLET, "tr", peak=0.18)
    running_mark(s, dark=True)

    y = Inches(1.14)
    y += eyebrow(s, M, y, d["eyebrow"], SKY) + Inches(0.2)
    y += headline(s, M, y, Inches(8.0), d["title"], ON_DARK, size=34)
    y += Inches(0.24)
    block(s, M, y, Inches(7.4), d["lead"], size=11.5, colour=ON_DARK_SOFT, leading=1.45)

    cols, gap = 4, Inches(0.46)
    col_w = Emu(int((CONTENT - gap * (cols - 1)) / cols))
    for i, (label, items) in enumerate(d["categories"]):
        col, row = i % cols, i // cols
        listing(s, M + col * (col_w + gap), Inches(2.84) + row * Inches(1.98),
                col_w, label, items, ACCENTS[i % len(ACCENTS)],
                head_colour=ON_DARK, item_colour=ON_DARK_SOFT, size=9,
                gap=Inches(0.038), leading=1.32)

    page_number(s, n, dark=True)
    notes(s, d["notes"])
    return s


def slide_dashboard(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False)
    running_mark(s, dark=False)

    col = Inches(4.24)
    y = Inches(1.14)
    y += eyebrow(s, M, y, d["eyebrow"], VIOLET) + Inches(0.2)
    y += headline(s, M, y, col, d["title"], INK, size=26)
    y += Inches(0.24)
    y += block(s, M, y, col, d["lead"], size=10.5, colour=INK_SOFT, leading=1.46)
    y += Inches(0.3)
    hairline(s, M, y, col, RULE)
    bullets(s, M, y + Inches(0.26), col, d["points"], size=10.2, gap=Inches(0.155))

    screenshot(s, d["shot"], Inches(5.6), Inches(1.14), Inches(6.87),
               crop=(0, 0, 0, 0.215), caption=d["caption"])

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_admissions(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False)
    running_mark(s, dark=False)

    col = Inches(5.4)
    y = Inches(1.14)
    y += eyebrow(s, M, y, d["eyebrow"], ROSE) + Inches(0.2)
    y += headline(s, M, y, col, d["title"], INK, size=29)
    y += Inches(0.26)
    y += block(s, M, y, col, d["lead"], size=10.8, colour=INK_SOFT, leading=1.46)
    y += Inches(0.34)
    def_grid(s, M, y, col, d["points"], cols=2, size_t=10.5, size_b=8.6,
             gap_x=Inches(0.34), gap_y=Inches(0.24))

    screenshot(s, d["shot"], Inches(6.86), Inches(1.12), Inches(5.61),
               crop=(0.157, 0.093, 0.005, 0.212))
    screenshot(s, d["shot2"], Inches(6.86), Inches(4.54), Inches(5.61),
               crop=(0.157, 0.118, 0.005, 0.53), caption=d["caption"])

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_records(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False)
    running_mark(s, dark=False)

    y = Inches(1.14)
    y += eyebrow(s, M, y, d["eyebrow"], SKY) + Inches(0.2)
    y += headline(s, M, y, Inches(6.2), d["title"], INK, size=31)
    y += Inches(0.26)
    block(s, M, y, Inches(5.4), d["lead"], size=11, colour=INK_SOFT, leading=1.46)

    gap = Inches(0.5)
    col_w = Emu(int((CONTENT - gap * 2) / 3))
    for i, (head, items) in enumerate(d["columns"]):
        listing(s, M + i * (col_w + gap), Inches(3.98), col_w, head, items,
                ACCENTS[i], size=9.5, gap=Inches(0.115))

    screenshot(s, d["shot"], Inches(7.62), Inches(1.1), Inches(4.86),
               crop=(0.157, 0.05, 0.005, 0.40), caption=d["caption"])

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_teachers(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False)
    running_mark(s, dark=False)

    col = Inches(4.5)
    y = Inches(1.14)
    y += eyebrow(s, M, y, d["eyebrow"], AMBER) + Inches(0.2)
    y += headline(s, M, y, col, d["title"], INK, size=31)
    y += Inches(0.28)
    block(s, M, y, col, d["lead"], size=11, colour=INK_SOFT, leading=1.46)

    screenshot(s, d["shot"], Inches(5.84), Inches(1.14), Inches(6.63),
               crop=(0.157, 0.05, 0.005, 0.435), caption=d["caption"])

    # Six points is more than a half column can hold beside a capture, so they
    # run as a band across the foot of the slide instead.
    def_grid(s, M, Inches(4.66), CONTENT, d["points"], cols=3,
             size_t=11, size_b=9, gap_x=Inches(0.52), gap_y=Inches(0.34))

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_engagement(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False, tone=PAPER_2)
    running_mark(s, dark=False)

    y = Inches(1.14)
    y += eyebrow(s, M, y, d["eyebrow"], MINT) + Inches(0.2)
    y += headline(s, M, y, Inches(6.0), d["title"], INK, size=31)
    y += Inches(0.26)
    block(s, M, y, Inches(5.2), d["lead"], size=11.5, colour=INK_SOFT, leading=1.46)

    card_y, card_h, col_w = Inches(3.62), Inches(3.02), Inches(3.42)
    for x, head, items, accent in [
        (M, d["left_head"], d["left"], VIOLET),
        (Inches(4.94), d["right_head"], d["right"], MINT),
    ]:
        card = rect(s, x - Inches(0.32), card_y, col_w + Inches(0.64), card_h,
                    fill=WHITE, line=RULE, radius=0.055)
        _shadow(card, blur=22, dist=7, alpha=0.07)
        cy = card_y + Inches(0.28)
        cy += block(s, x, cy, col_w, head.upper(), size=9.5, colour=accent,
                    font=DISPLAY, spc=130, leading=1.25)
        bullets(s, x, cy + Inches(0.18), col_w, items, size=9.5, dot=accent,
                gap=Inches(0.1), leading=1.3)

    # A phone, because that is the only device a parent will ever open it on,
    # with the feedback figures as a detail strip beneath it.
    screenshot(s, d["phone"], Inches(9.86), Inches(1.12), Inches(1.92),
               crop=(0, 0, 0, 0.30))
    screenshot(s, d["shot"], Inches(9.16), Inches(4.3), Inches(3.32),
               crop=(0.157, 0.085, 0.015, 0.70))
    block(s, Inches(9.16), Inches(5.06), Inches(3.32), d["caption"], size=9.5,
          colour=INK_FAINT, leading=1.34)

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_operations(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False)
    running_mark(s, dark=False)

    y = Inches(1.14)
    y += eyebrow(s, M, y, d["eyebrow"], VIOLET) + Inches(0.2)
    y += headline(s, M, y, Inches(7.0), d["title"], INK, size=29)
    y += Inches(0.24)
    block(s, M, y, Inches(5.4), d["lead"], size=11, colour=INK_SOFT, leading=1.44)

    gap = Inches(0.5)
    col_w = Emu(int((CONTENT - gap * 2) / 3))
    for i, (head, items) in enumerate(d["columns"]):
        listing(s, M + i * (col_w + gap), Inches(4.0), col_w, head, items,
                ACCENTS[i], size=9.5, gap=Inches(0.115))

    screenshot(s, d["shot"], Inches(8.06), Inches(1.1), Inches(4.42),
               crop=(0.157, 0.05, 0.005, 0.325), caption=d["caption"])

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_assistant(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=True)
    wash(s, VIOLET, "bl", peak=0.22)
    running_mark(s, dark=True)

    col = Inches(4.16)
    y = Inches(1.12)
    y += eyebrow(s, M, y, d["eyebrow"], VIOLET_SOFT) + Inches(0.2)
    y += headline(s, M, y, col, d["title"], ON_DARK, size=27)
    y += Inches(0.24)
    y += block(s, M, y, col, d["lead"], size=10.5, colour=ON_DARK_SOFT, leading=1.44)
    y += Inches(0.3)

    for q in d["questions"]:
        chip = rect(s, M, y, col, Inches(0.34), fill=RGBColor(0x16, 0x1C, 0x33),
                    line=RULE_DARK, radius=0.3)
        _no_shadow(chip)
        block(s, M + Inches(0.22), y + Inches(0.08), col - Inches(0.4), q,
              size=10, colour=ON_DARK, leading=1.2)
        y += Inches(0.42)

    screenshot(s, d["shot"], Inches(5.44), Inches(1.12), Inches(7.03),
               crop=(0, 0, 0, 0.20), caption=d["caption"], caption_colour=ON_DARK_FAINT)

    # The guarantees run as a band beneath, where they read as the terms of the
    # thing rather than as a footnote to it.
    band_y = max(Inches(5.76), y + Inches(0.22))
    hairline(s, M, band_y, CONTENT, RULE_DARK)
    gap = Inches(0.44)
    cell = Emu(int((CONTENT - gap * 3) / 4))
    for i, (title, body) in enumerate(d["rules"]):
        cx = M + i * (cell + gap)
        ty = band_y + Inches(0.24)
        ty += block(s, cx, ty, cell, title, size=10, colour=ACCENTS[i % len(ACCENTS)],
                    font=DISPLAY, leading=1.25)
        block(s, cx, ty + Inches(0.06), cell, body, size=8.6, colour=ON_DARK_SOFT,
              leading=1.36)

    page_number(s, n, dark=True)
    notes(s, d["notes"])
    return s


def slide_impact(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=False)
    running_mark(s, dark=False)

    col = Inches(5.6)
    y = Inches(1.14)
    y += eyebrow(s, M, y, d["eyebrow"], MINT) + Inches(0.2)
    y += headline(s, M, y, col, d["title"], INK, size=30)
    y += Inches(0.3)
    y += block(s, M, y, col, d["impact_head"].upper(), size=9.5, colour=INK_FAINT,
               font=DISPLAY, spc=140)
    def_grid(s, M, y + Inches(0.22), col, d["impact"], cols=2, size_t=10.5,
             size_b=8.6, gap_x=Inches(0.42), gap_y=Inches(0.3))

    # The implementation sequence, as a numbered rail on its own ground.
    x = Inches(7.3)
    rail_w = W - M - x
    card = rect(s, x - Inches(0.44), Inches(1.1), rail_w + Inches(0.88), Inches(5.34),
                fill=PAPER_2, radius=0.032)
    _no_shadow(card)
    cy = Inches(1.46)
    cy += block(s, x, cy, rail_w, d["journey_head"].upper(), size=9.5, colour=VIOLET,
                font=DISPLAY, spc=140)
    cy += Inches(0.2)
    for i, (step, title, when, body) in enumerate(d["journey"]):
        block(s, x, cy, Inches(0.42), step, size=10.5,
              colour=ACCENTS[i % len(ACCENTS)], font=DISPLAY)
        block(s, x + Inches(0.46), cy - Inches(0.012), Inches(2.6), title, size=11.5,
              colour=INK, font=DISPLAY)
        block(s, x + Inches(2.4), cy, rail_w - Inches(2.4), when, size=8.8,
              colour=INK_FAINT, align=PP_ALIGN.RIGHT)
        cy += Inches(0.25)
        cy += block(s, x + Inches(0.46), cy, rail_w - Inches(0.46), body, size=8.6,
                    colour=INK_SOFT, leading=1.32)
        cy += Inches(0.15)

    page_number(s, n, dark=False)
    notes(s, d["notes"])
    return s


def slide_cta(prs, d, n):
    s = new_slide(prs)
    ground(s, dark=True, tone=INK_DEEP)
    wash(s, VIOLET, "tr", peak=0.42)
    wash(s, AMBER, "br", peak=0.13)

    lockup(s, M, Inches(0.66), Inches(0.5), dark=True)

    y = Inches(2.04)
    y += headline(s, M, y, Inches(6.4), d["title"], ON_DARK, size=38, leading=1.1)
    y += Inches(0.32)
    y += block(s, M, y, Inches(5.4), d["lead"], size=12.5, colour=ON_DARK_SOFT,
               leading=1.5)

    y += Inches(0.44)
    btn = rect(s, M, y, Inches(2.16), Inches(0.6), fill=WHITE, radius=0.5)
    _shadow(btn, blur=26, dist=9, alpha=0.28)
    block(s, M, y + Inches(0.17), Inches(2.16), "Book a demo", size=13.5, colour=INK,
          font=DISPLAY, align=PP_ALIGN.CENTER)
    block(s, M + Inches(2.44), y + Inches(0.19), Inches(3.4), C.SALES_EMAIL, size=12.5,
          colour=ON_DARK)

    hairline(s, M, Inches(6.04), Inches(5.4), RULE_DARK)
    block(s, M, Inches(6.26), Inches(5.4), d["kicker"], size=10, colour=ON_DARK_SOFT,
          leading=1.42)

    x = Inches(7.66)
    inner = W - M - x
    cy = Inches(2.1)
    cy += block(s, x, cy, inner, d["involves_head"].upper(), size=9.5,
                colour=AMBER, font=DISPLAY, spc=140)
    cy += Inches(0.2)
    for item in d["involves"]:
        hairline(s, x, cy, inner, RULE_DARK)
        cy += Inches(0.22)
        cy += block(s, x, cy, inner, item, size=11.5, colour=ON_DARK_SOFT, leading=1.42)
        cy += Inches(0.26)

    block(s, x, Inches(6.26), inner, C.TAGLINE, size=11, colour=AMBER,
          font=DISPLAY, spc=60)

    page_number(s, n, dark=True)
    notes(s, d["notes"])
    return s


# ---------------------------------------------------------------------- main

def build():
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H

    slide_cover(prs, C.COVER)
    slide_challenge(prs, C.CHALLENGE, 2)
    slide_connected(prs, C.CONNECTED, 3)
    slide_what(prs, C.WHAT, 4)
    slide_value(prs, C.VALUE, 5)
    slide_platform(prs, C.PLATFORM, 6)
    slide_dashboard(prs, C.DASHBOARD, 7)
    slide_admissions(prs, C.ADMISSIONS, 8)
    slide_records(prs, C.RECORDS, 9)
    slide_teachers(prs, C.TEACHERS, 10)
    slide_engagement(prs, C.ENGAGEMENT, 11)
    slide_operations(prs, C.OPERATIONS, 12)
    slide_assistant(prs, C.ASSISTANT, 13)
    slide_impact(prs, C.IMPACT, 14)
    slide_cta(prs, C.CTA, 15)

    core = prs.core_properties
    core.title = "MyCampusView — one operating system for your entire school"
    core.author = "MyCampusView"
    core.subject = ("School management platform: student information system, admission CRM "
                    "and school ERP on one database")
    core.category = "Sales"
    core.comments = ("Built by docs/sales-deck/build_deck.py. Copy lives in deck_content.py; "
                     "screenshots are real captures of the running application.")

    prs.save(OUT)
    print(f"{OUT.name}  ·  {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    build()
